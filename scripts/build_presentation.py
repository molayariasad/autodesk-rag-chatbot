"""
Generate the 20-minute panel presentation for the Autodesk GET-ML interview.

Covers:
1. Approach to the problem, data, and architecture
2. Production system plan (if PoC → full product)
3. Reflections: preprocessing, doc selection, RAG architecture, evaluation
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# Color Palette — Ocean Gradient (matches Autodesk brand feel)
# ============================================================
NAVY = RGBColor(0x06, 0x5A, 0x82)
TEAL = RGBColor(0x1C, 0x72, 0x93)
DARK = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
MID_GRAY = RGBColor(0x64, 0x74, 0x8B)
ACCENT = RGBColor(0x02, 0xC3, 0x9A)
SOFT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)

FONT_TITLE = "Georgia"
FONT_BODY = "Calibri"


def add_dark_slide(prs, title_text, subtitle_text=None):
    """Dark background slide for section openers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_TITLE
    p.alignment = PP_ALIGN.LEFT

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(18)
        p2.font.color.rgb = LIGHT_GRAY
        p2.font.name = FONT_BODY
        p2.space_before = Pt(16)

    # Accent bar
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.25), Inches(1.5), Inches(0.06)
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = ACCENT
    slide.shapes[-1].line.fill.background()

    return slide


def add_content_slide(prs, title_text, bullet_points, note_text=None):
    """Light background content slide with bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    # Left accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(5.625)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.font.name = FONT_TITLE

    # Bullet content
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(3.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()

        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.font.name = FONT_BODY
        p.space_after = Pt(10)
        p.level = 0

        # Manual bullet
        p.text = "▸  " + point

    # Optional note at bottom
    if note_text:
        txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(8.4), Inches(0.5))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = note_text
        p3.font.size = Pt(11)
        p3.font.color.rgb = MID_GRAY
        p3.font.italic = True
        p3.font.name = FONT_BODY

    return slide


def add_two_column_slide(prs, title_text, left_items, right_items, left_header="", right_header=""):
    """Two-column comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    # Left accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(5.625)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.font.name = FONT_TITLE

    # Left column
    def add_column(x, header, items):
        if header:
            hbox = slide.shapes.add_textbox(Inches(x), Inches(1.2), Inches(4.2), Inches(0.5))
            hp = hbox.text_frame.paragraphs[0]
            hp.text = header
            hp.font.size = Pt(18)
            hp.font.bold = True
            hp.font.color.rgb = NAVY
            hp.font.name = FONT_BODY

        cbox = slide.shapes.add_textbox(Inches(x), Inches(1.8), Inches(4.2), Inches(3.5))
        ctf = cbox.text_frame
        ctf.word_wrap = True
        for i, item in enumerate(items):
            cp = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
            cp.text = "▸  " + item
            cp.font.size = Pt(14)
            cp.font.color.rgb = DARK
            cp.font.name = FONT_BODY
            cp.space_after = Pt(8)

    add_column(0.7, left_header, left_items)
    add_column(5.2, right_header, right_items)

    # Divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.95), Inches(1.3), Inches(0.02), Inches(3.8)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    line.line.fill.background()

    return slide


def add_stat_slide(prs, title_text, stats):
    """Big numbers slide for key metrics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(8.4), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_TITLE

    # Stats in a row
    n = len(stats)
    col_w = 8.4 / n
    for i, (value, label) in enumerate(stats):
        x = 0.8 + i * col_w

        # Value
        vbox = slide.shapes.add_textbox(Inches(x), Inches(1.8), Inches(col_w - 0.2), Inches(1.2))
        vp = vbox.text_frame.paragraphs[0]
        vp.text = value
        vp.font.size = Pt(32)
        vp.font.bold = True
        vp.font.color.rgb = ACCENT
        vp.font.name = FONT_TITLE
        vp.alignment = PP_ALIGN.CENTER

        # Label
        lbox = slide.shapes.add_textbox(Inches(x), Inches(3.0), Inches(col_w - 0.2), Inches(0.8))
        lp = lbox.text_frame.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(14)
        lp.font.color.rgb = LIGHT_GRAY
        lp.font.name = FONT_BODY
        lp.alignment = PP_ALIGN.CENTER

    return slide


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # ============================================================
    # SLIDE 1: Title
    # ============================================================
    slide = add_dark_slide(
        prs,
        "Autodesk RAG Chatbot",
        "Senior ML Engineer Take-Home Assessment\nAsad Molayari"
    )

    # ============================================================
    # SLIDE 2: Problem Statement
    # ============================================================
    add_content_slide(prs, "Problem Statement", [
        "Build a RAG chatbot over 1,218 noisy HTML pages from Autodesk's website",
        "Two retrieval modes: corpus-only and blended (corpus + web search)",
        "Evaluate retrieval quality, generation faithfulness, and system performance",
        "Demonstrate Senior-level system design, data engineering, and ML skills",
    ], "Key challenge: HTML pages are highly heterogeneous — JS shells, marketing pages, help docs, e-commerce")

    # ============================================================
    # SLIDE 3: Architecture Overview
    # ============================================================
    add_content_slide(prs, "System Architecture", [
        "Streamlit UI → FastAPI → RAG Chain → Ollama (Mistral 7B)",
        "Hybrid Retrieval: Vector Search (BGE) + BM25 → RRF Fusion → Cross-Encoder Reranking",
        "Vector Store: ChromaDB (embedded, persistent, no separate server)",
        "Two modes: corpus-only and blended (DuckDuckGo web augmentation)",
        "Fully containerized: docker-compose up spins up Ollama + FastAPI + Streamlit",
    ])

    # ============================================================
    # SLIDE 4: Tech Stack
    # ============================================================
    add_two_column_slide(
        prs,
        "Tech Stack Decisions",
        [
            "Mistral 7B Instruct — strong instruction-following, runs locally",
            "BGE-small-en-v1.5 — top MTEB scores at 33M params",
            "ChromaDB — embedded, zero-config, ideal for PoC",
            "BM25 (rank_bm25) — catches exact product names",
            "Cross-Encoder reranking — precise relevance scoring",
        ],
        [
            "FastAPI — async, auto-docs, Pydantic validation",
            "Streamlit — rapid chat UI with built-in components",
            "UV — 10-100x faster Python package management",
            "Docker + docker-compose — one-command deployment",
            "DuckDuckGo Search — no API key, privacy-friendly",
        ],
        "ML & Retrieval",
        "Infrastructure",
    )

    # ============================================================
    # SLIDE 5: Data Challenge
    # ============================================================
    add_stat_slide(prs, "The Data Challenge", [
        ("1,218", "HTML Files"),
        ("5KB–3MB", "File Size Range"),
        ("~40%", "Blank/JS-Only Pages"),
        ("0–37", "Scripts Per Page"),
    ])

    # ============================================================
    # SLIDE 6: HTML Preprocessing
    # ============================================================
    add_content_slide(prs, "Text Preprocessing: Multi-Layer Extraction", [
        "Layer 1 — Trafilatura: Best for article/marketing content (Forma blog, product pages)",
        "Layer 2 — BeautifulSoup: Targets main/article divs (help documentation pages)",
        "Layer 3 — html2text: Fallback for remaining edge cases",
        "Table preservation: HTML tables → Markdown before boilerplate removal",
        "Quality gate: Skip pages with < 100 chars of meaningful text",
    ], "Each layer handles a different HTML archetype — no single parser works for all Autodesk page types")

    # ============================================================
    # SLIDE 7: Document Selection
    # ============================================================
    add_content_slide(prs, "Document Selection Strategy", [
        "Ingest all documents that pass quality filtering (> 100 chars meaningful text)",
        "Blank JS shells and loader pages are automatically excluded",
        "Metadata enrichment: product name, page type, title extracted from <meta> tags",
        "Chunking: 512 tokens, 64 overlap, with metadata prefix for context",
        "Table chunks preserved separately to maintain structural integrity",
    ], "Rationale: With 1,218 files, storage is not a bottleneck — better to have coverage than miss relevant docs")

    # ============================================================
    # SLIDE 8: Retrieval Strategy
    # ============================================================
    add_content_slide(prs, "Hybrid Retrieval Pipeline", [
        "Stage 1 — Vector Search: BGE embeddings + ChromaDB cosine similarity (top-10)",
        "Stage 2 — BM25 Search: Keyword matching catches exact product names (top-10)",
        "Stage 3 — Reciprocal Rank Fusion: Merges both lists (weight: 0.7 vector, 0.3 BM25)",
        "Stage 4 — Cross-Encoder Reranking: Joint (query, passage) scoring for precision",
        "Stage 5 — Top-5 context passed to LLM with citation instructions",
    ], "RRF avoids the score normalization problem — proven in Cormack et al., 2009")

    # ============================================================
    # SLIDE 9: Why Hybrid?
    # ============================================================
    add_two_column_slide(
        prs,
        "Why Hybrid Search Matters",
        [
            "Semantic: 'What 3D modeling tool?' matches 'create 3D objects'",
            "Cannot match: 'AutoCAD LT 2024' exactly",
            "Great for natural language questions",
            "Embedding captures intent, not tokens",
        ],
        [
            "Keyword: 'AutoCAD LT' matches exactly in BM25",
            "Struggles with: paraphrased questions",
            "Great for product names, versions, specs",
            "BM25 catches what embeddings miss",
        ],
        "Vector Search (Semantic)",
        "BM25 (Keyword)",
    )

    # ============================================================
    # SLIDE 10: Hallucination Mitigation
    # ============================================================
    add_content_slide(prs, "Hallucination Mitigation", [
        "Grounded prompting: 'ONLY answer from provided context — say I don't know if unsure'",
        "Citation enforcement: Prompt instructs model to cite [Source: Title] for every claim",
        "Cross-encoder reranking ensures most relevant (not just similar) passages reach LLM",
        "Heuristic detection: Check if specific claims (prices, versions) appear in context",
        "Graceful refusal: System explicitly declines off-topic questions",
    ])

    # ============================================================
    # SLIDE 11: Conversation Support
    # ============================================================
    add_content_slide(prs, "Conversation & Follow-Up Support", [
        "Chat history (last 3 turns) included in every prompt for coreference resolution",
        "Example: 'Tell me about Maya' → 'What about its pricing?' → resolves 'its' to Maya",
        "Streamlit UI persists full conversation with source inspection panel",
        "Thumbs up/down feedback collected per response for future fine-tuning",
        "Source transparency: Users can inspect which documents informed each answer",
    ])

    # ============================================================
    # SLIDE 12: Evaluation Framework
    # ============================================================
    add_content_slide(prs, "Evaluation Framework", [
        "Retrieval metrics: Hit Rate, Keyword Overlap, Product Match Rate",
        "Generation metrics: Answer Relevance, Citation Rate, Refusal Accuracy",
        "System metrics: Latency, Hallucination Flags, Error Count",
        "10 test questions: 5 from assignment + 5 additional (pricing, comparison, irrelevant)",
        "All results saved as structured JSON — reproducible, extensible",
    ], "Deterministic metrics first (reproducible) → LLM-as-judge in production (deeper quality)")

    # ============================================================
    # SLIDE 13: Evaluation Validity
    # ============================================================
    add_content_slide(prs, "Evaluation: What & How We Measure", [
        "WHAT: Retrieval quality (right docs?), Generation faithfulness (grounded?), Coverage (keywords?)",
        "HOW: Keyword overlap for relevance, regex for citations, claim extraction for hallucination",
        "VALIDITY: Deterministic metrics are reproducible but limited in semantic depth",
        "Honest limitation: Keyword overlap ≠ true 'correctness' — it's a practical proxy",
        "Next step: LLM-as-judge (GPT-4/Claude) for deeper semantic evaluation",
    ])

    # ============================================================
    # SLIDE 14: Corpus vs Blended
    # ============================================================
    add_two_column_slide(
        prs,
        "Corpus-Only vs. Blended Mode",
        [
            "Answers strictly from Autodesk HTML docs",
            "Higher precision: all info is authoritative",
            "Limited recall: some questions have incomplete answers",
            "No external data leakage risk",
            "Best for: product specs, help documentation",
        ],
        [
            "Augments with DuckDuckGo web results",
            "Higher recall: fills gaps in corpus",
            "Lower precision: web info may be outdated",
            "Source attribution distinguishes internal vs web",
            "Best for: latest releases, pricing, comparisons",
        ],
        "Corpus-Only Mode",
        "Blended Mode",
    )

    # ============================================================
    # SLIDE 15: Production Roadmap
    # ============================================================
    add_dark_slide(
        prs,
        "From PoC to Production",
        "What would a full-featured product look like?"
    )

    # ============================================================
    # SLIDE 16: Week 1-2
    # ============================================================
    add_two_column_slide(
        prs,
        "Production Plan: Weeks 1–2",
        [
            "Headless browser crawl (Playwright) for JS-rendered pages",
            "Structured metadata extraction pipeline",
            "PostgreSQL for conversations and feedback storage",
            "Migrate ChromaDB → Qdrant (production-grade)",
            "CI/CD pipeline with automated testing",
        ],
        [
            "Fine-tune BGE embeddings on Autodesk domain data",
            "Query understanding: intent classification + entity extraction",
            "Two-stage retrieval: document-level then chunk-level",
            "A/B test chunk sizes and overlap ratios",
            "Expand evaluation with LLM-as-judge scoring",
        ],
        "Week 1: Data & Infrastructure",
        "Week 2: Retrieval & Model Quality",
    )

    # ============================================================
    # SLIDE 17: Week 3-4
    # ============================================================
    add_two_column_slide(
        prs,
        "Production Plan: Weeks 3–4",
        [
            "LangGraph multi-agent architecture",
            "Router agent: classify intent (product, pricing, support)",
            "Tool-use for structured queries (pricing API, catalog)",
            "Streaming responses for better UX",
            "RLHF from collected human feedback",
        ],
        [
            "Kubernetes deployment with auto-scaling",
            "Monitoring: latency P99, error rates, hallucination tracking",
            "User analytics dashboard (popular queries, failure patterns)",
            "Automated regression testing on config changes",
            "Evaluation gates in CI/CD pipeline",
        ],
        "Week 3: Agent & Multi-Turn",
        "Week 4: Deploy & Monitor",
    )

    # ============================================================
    # SLIDE 18: Key Reflections
    # ============================================================
    add_content_slide(prs, "Key Reflections", [
        "Data quality is the #1 lever: Better HTML parsing → better chunks → better retrieval → better answers",
        "Hybrid search is not optional: BM25 catches 'AutoCAD LT' that vector search misses",
        "Reranking is high-ROI: Cross-encoder dramatically improves top-k precision for small cost",
        "Evaluation is iterative: Started with keyword overlap, added hallucination detection, next: LLM-as-judge",
        "Open-source LLMs are production-viable: Mistral 7B gives solid RAG performance locally",
    ])

    # ============================================================
    # SLIDE 19: Summary
    # ============================================================
    add_stat_slide(prs, "System Summary", [
        ("3,000+", "Lines of Code"),
        ("3-Layer", "HTML Parsing"),
        ("4-Stage", "Hybrid Retrieval"),
        ("10", "Eval Questions"),
    ])

    # ============================================================
    # SLIDE 20: Thank You / Q&A
    # ============================================================
    slide = add_dark_slide(
        prs,
        "Thank You",
        "Questions & Discussion\n\nAsad Molayari  |  amolayari.s@gmail.com"
    )

    return prs


if __name__ == "__main__":
    prs = build_presentation()
    output_path = "/home/claude/presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
